import os
from pptx import Presentation
from pptx.util import Inches

def build_pptx():
    prs = Presentation()
    # 设置 16:9 比例 (宽 13.333 英寸，高 7.5 英寸)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 幻灯片布局，我们使用空白布局（空白页）
    blank_slide_layout = prs.slide_layouts[6]
    
    img_dir = '/Users/gzm/Downloads/手机剪辑课程/_pdf_frames'
    output_path = '/Users/gzm/Downloads/手机剪辑课程/手机剪辑零基础入门课_图片版.pptx'
    
    # 循环 52 页并塞入图片
    for i in range(1, 53):
        img_name = f'slide-{i:02d}.png'
        img_path = os.path.join(img_dir, img_name)
        
        if not os.path.exists(img_path):
            print(f"Warning: {img_path} 不存在，已跳过。")
            continue
            
        # 添加新的一页 slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 将图片添加并铺满整个幻灯片
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
    prs.save(output_path)
    print(f"PPTX 课件成功导出并保存至：{output_path}")

if __name__ == '__main__':
    build_pptx()
